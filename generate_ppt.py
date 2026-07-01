import collections
import collections.abc
import os
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx is not installed. Please run: pip install python-pptx")
    exit(1)

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # Widescreen 16:9
prs.slide_height = Inches(7.5)

# Slide data
slides_data = [
    {
        "title": "Dataset Curation",
        "bullets": [
            "Multi-Terrain Representation: Thar Desert, Delhi Urban, Kutch Coastal, Ladakh Mountain, Telangana Agriculture.",
            "Sensor Radiometric Calibration: DN converted to Celsius LST; visible bands normalized to [0,1] range.",
            "Sliding-Window Patching: Co-registered 128x128 pixel paired spectral patches.",
            "Anti-Leakage Geographic Filter: Spatial Block Splitting (80/20) prevents coordinate overlap."
        ]
    },
    {
        "title": "Training Strategy",
        "bullets": [
            "Decoupled Optimization Loops: Spatial upscaling (SwinIR) separated from spectral translation (cGAN).",
            "Heteroscedastic Uncertainty Loss: Mathematical log-likelihood scales loss by predicted variance.",
            "Loss Attenuation: Automatically discounts penalty on ambiguous pixels, preventing gradient oscillations.",
            "Geographically Blind Validation: Model evaluated on geographically isolated validation splits at every epoch."
        ]
    },
    {
        "title": "Technical Credibility",
        "bullets": [
            "Decoupled Spatial-Spectral Pipeline: Lowers the mathematical search space, improving convergence.",
            "High-Fidelity Metric Benchmarks: Authentic validation yields PSNR >40 dB and SSIM >0.97.",
            "No-Reference Quality Metrics: Details, colorfulness, and edge sharpness computed for blind uploads.",
            "Rigorous Validation Split: Spatial block partitioning prevents training coordinate contamination."
        ]
    },
    {
        "title": "Novelty & USP",
        "bullets": [
            "Self-Supervised Uncertainty Estimation: Network outputs RGB prediction alongside a pixel-wise variance map.",
            "Uncertainty-Guided Physics Fusion: Blends neural outputs with a remote-sensing Ts-NDVI physical model.",
            "Guided Structure Detail Transfer: High-res thermal output acts as a structural guide to prevent blurred boundaries.",
            "Multi-Scale Laplacian Pyramid Injection: Direct detail transfer from thermal bands preserves high-frequency edges."
        ]
    }
]

# Style settings
dark_blue = RGBColor(11, 19, 43)
light_cyan = RGBColor(100, 223, 223)
white = RGBColor(255, 255, 255)
gray = RGBColor(220, 225, 230)

blank_layout = prs.slide_layouts[6]

for data in slides_data:
    slide = prs.slides.add_slide(blank_layout)
    
    # Set dark background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = dark_blue
    
    # 1. Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.0))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = data["title"]
    p_title.font.name = "Arial"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = light_cyan
    
    # 2. Slide Container Card (Visual Box)
    shape = slide.shapes.add_shape(
        1, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.7)  # 1 = MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(20, 30, 55)
    shape.line.color.rgb = RGBColor(100, 223, 223)
    shape.line.width = Pt(1.5)
    
    # 3. Content Text Frame inside Card
    tf_content = shape.text_frame
    tf_content.word_wrap = True
    tf_content.margin_left = Inches(0.4)
    tf_content.margin_right = Inches(0.4)
    tf_content.margin_top = Inches(0.4)
    
    for i, bullet in enumerate(data["bullets"]):
        p = tf_content.add_paragraph() if i > 0 else tf_content.paragraphs[0]
        p.text = "▪  " + bullet
        p.space_after = Pt(20)
        p.font.name = "Arial"
        p.font.size = Pt(20)
        p.font.color.rgb = gray

# Save the presentation
output_name = "isro_bah_pitch.pptx"
prs.save(output_name)
print(f"Presentation generated successfully: '{output_name}'")
